from pptx import Presentation
import sys

def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        
        print(f"=" * 80)
        print(f"文件: {file_path}")
        print(f"幻灯片数量: {len(prs.slides)}")
        print(f"=" * 80)
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n{'='*60}")
            print(f"幻灯片 {i}")
            print(f"{'='*60}")
            
            # 获取幻灯片布局名称
            if slide.slide_layout:
                print(f"布局: {slide.slide_layout.name}")
            
            # 遍历所有形状
            for shape in slide.shapes:
                # 处理文本框和占位符
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"\n[文本框] {shape.shape_type}:")
                    print(f"{shape.text}")
                
                # 处理表格
                if shape.has_table:
                    print(f"\n[表格]:")
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        print(f"  行{row_idx}: {' | '.join(row_text)}")
        
        return True
    except Exception as e:
        print(f"错误: {e}")
        return False

if __name__ == "__main__":
    file_path = r"d:\little project\genecloud\第06章 数据库应用实践.pptx"
    read_pptx(file_path)